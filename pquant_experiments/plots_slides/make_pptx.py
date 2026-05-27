#!/usr/bin/env python3
"""
make_pptx.py — Editable PowerPoint slide deck for ParT wf scan results.

Outputs:
  scan_slides.pptx  — 5 slides: accuracy plot, AUC plot, overview table,
                       per-class AUC table, QCD rejection table.

Editability:
  - Titles, subtitles, footnotes  → native PowerPoint text boxes
  - Tables                        → native PowerPoint tables (fully editable)
  - Plots                         → embedded high-res PNG images

Usage:
    cd /eos/home-m/majorgen/particle_transformer
    python3 pquant_experiments/make_pptx.py
    python3 pquant_experiments/make_pptx.py --output scan_slides.pptx
"""

import argparse
import io
import re

import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import wandb

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree

mpl.rcParams["font.family"] = "sans-serif"

# ── Project constants ─────────────────────────────────────────────────────────
ENTITY       = "martina-jorgensen-cern"
PROJECT      = "par-t-quant"
BASELINE_ACC = 86.05
BASELINE_RW  = 1.0

CLASSES        = ['QCD', 'Hbb', 'Hcc', 'Hgg', 'H4q', 'Hqql', 'Zqq', 'Wqq', 'Tbqq', 'Tbl']
SIGNAL_CLASSES = ['Hbb', 'Hcc', 'Hgg', 'H4q', 'Hqql', 'Zqq', 'Wqq', 'Tbqq', 'Tbl']
TPR_TARGETS    = {
    'Hbb': 0.50, 'Hcc': 0.50, 'Hgg': 0.50, 'H4q': 0.50,
    'Hqql': 0.99, 'Zqq': 0.50, 'Wqq': 0.50, 'Tbqq': 0.50, 'Tbl': 0.995,
}

BASELINE = {
    "wf": "FP", "run_name": "full-precision",
    "acc_%": 86.051, "acc_is_test": True, "remaining_w": 1.0,
    "test_auc_ovo": 0.987720,
    "auc_QCD": 0.978251, "auc_Hbb": 0.997134, "auc_Hcc": 0.987901,
    "auc_Hgg": 0.980071, "auc_H4q": 0.988796, "auc_Hqql": 0.999567,
    "auc_Zqq": 0.966947, "auc_Wqq": 0.980131, "auc_Tbqq": 0.998537,
    "auc_Tbl": 0.999865,
    "rej_Hbb": 10638, "rej_Hcc": 4149, "rej_Hgg": 123, "rej_H4q": 1867,
    "rej_Hqql": 5420, "rej_Zqq": 402, "rej_Wqq": 543,
    "rej_Tbqq": 32258, "rej_Tbl": 16260,
}

# ── Colours (RGBColor for pptx, hex for matplotlib) ──────────────────────────
def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

C_BLUE     = "#0073E6"
C_TEXT     = "#1a2233"
C_SUBTEXT  = "#4a5568"
C_GREY     = "#9aacbb"
C_GRID     = "#e8edf2"
C_FRAME    = "#d0d7df"
C_HDR_BG   = "#dbeafe"
C_HDR_TEXT = "#1e40af"
C_ROW_A    = "#ffffff"
C_ROW_B    = "#f1f5f9"
C_BL_ROW   = "#eff6ff"   # baseline row tint
C_BORDER   = "#e2e8f0"
C_WHITE    = "#ffffff"

# Slide dimensions: 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.55)   # ~1.4 cm
STRIPE_W = Inches(0.12)
HEADER_H = Inches(1.45)  # height reserved for title + subtitle
PLOT_SUBTITLE = "Data compressed to 16 bits (f=11, i=4, k=1)"


# ── Helpers ───────────────────────────────────────────────────────────────────

def get_wf(run):
    wf = run.config.get("WEIGHT_F_BITS") or run.config.get("weight_f_bits")
    if wf is not None:
        return int(wf)
    m = re.search(r"wf(\d+)", run.name or "")
    if m:
        return int(m.group(1))
    return None


def _pct(v):
    if v is None:
        return None
    return float(v) * 100 if float(v) <= 1.0 else float(v)


def _nan(v):
    return v is None or (isinstance(v, float) and np.isnan(v))


def fmt_acc(v): return f"{v:.3f}%" if not _nan(v) else "—"
def fmt_auc(v): return f"{v:.5f}"  if not _nan(v) else "—"
def fmt_wt(v):
    if _nan(v): return "—"
    p = float(v) * 100 if float(v) <= 1.0 else float(v)
    return f"{p:.2f}%"
def fmt_rej(v):
    if _nan(v): return "—"
    return "∞" if v == float("inf") else f"{v:.3f}"


# ── Data fetching ─────────────────────────────────────────────────────────────

def fetch_full_data():
    api  = wandb.Api()
    runs = api.runs(f"{ENTITY}/{PROJECT}")
    records = []
    for run in runs:
        if run.state not in ("finished", "running", "crashed"):
            continue
        wf = get_wf(run)
        if wf is None or wf == 14:
            continue
        s        = run.summary
        test_acc = s.get("test/acc")
        auc_ovo  = s.get("test/auc_ovo_macro")
        eval_acc = _pct(s.get("eval/acc_epoch"))
        rem_w    = s.get("eval/remaining_weights")
        acc         = test_acc if test_acc is not None else eval_acc
        acc_is_test = test_acc is not None
        if _nan(acc) and _nan(auc_ovo):
            continue
        row = {"wf": wf, "run_name": run.name, "acc_%": acc,
               "acc_is_test": acc_is_test, "remaining_w": rem_w,
               "test_auc_ovo": auc_ovo, "state": run.state}
        for cls in CLASSES:
            row[f"auc_{cls}"] = s.get(f"test/auc_{cls}")
        for cls in SIGNAL_CLASSES:
            lbl = int(TPR_TARGETS[cls] * 100)
            row[f"rej_{cls}"] = s.get(f"test/rejection_{cls}_at{lbl}eff")
        records.append(row)
        src = "test" if acc_is_test else "eval/last"
        print(f"  wf={wf:2d}  acc={fmt_acc(acc)}({src})  auc={fmt_auc(auc_ovo)}")

    if not records:
        return pd.DataFrame()
    df = pd.DataFrame(records)
    df = df.sort_values(["wf", "acc_%"], na_position="first")
    df = df.drop_duplicates(subset="wf", keep="last")
    df = df.sort_values("wf").reset_index(drop=True)
    return df


# ── pptx slide utilities ──────────────────────────────────────────────────────

def new_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank layout
    slide  = prs.slides.add_slide(layout)
    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(C_WHITE)
    return slide


def add_stripe(slide):
    """Thin blue vertical stripe on the left edge."""
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), STRIPE_W, SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(C_BLUE)
    shape.line.fill.background()   # no border


def add_title_block(slide, title, subtitle=None):
    """Add editable title + optional subtitle text boxes."""
    left   = STRIPE_W + MARGIN
    width  = SLIDE_W - STRIPE_W - MARGIN - Inches(0.3)

    # Title
    tb = slide.shapes.add_textbox(left, MARGIN * 0.6, width, Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.bold  = True
    run.font.size  = Pt(24)
    run.font.color.rgb = rgb(C_TEXT)

    # Subtitle
    if subtitle:
        tb2 = slide.shapes.add_textbox(left, MARGIN * 0.6 + Inches(0.72),
                                        width, Inches(0.40))
        tf2 = tb2.text_frame
        tf2.word_wrap = False
        p2  = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size  = Pt(13)
        run2.font.color.rgb = rgb(C_BLUE)


def add_footnote(slide, text):
    """Small italic footnote at the bottom of the slide."""
    tb = slide.shapes.add_textbox(
        STRIPE_W + MARGIN,
        SLIDE_H - Inches(0.42),
        SLIDE_W - STRIPE_W - 2 * MARGIN,
        Inches(0.38)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(8)
    run.font.italic = True
    run.font.color.rgb = rgb(C_GREY)


def embed_plot_image(slide, fig, left, top, width, height):
    """Render a matplotlib figure to PNG in-memory and embed it."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width, height)
    plt.close(fig)


# ── pptx table builder ────────────────────────────────────────────────────────

def set_cell_color(cell, hex_color):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", hex_color.lstrip("#"))


def set_cell_border(cell, hex_color, width_pt=0.5):
    """Add thin border on all sides of a cell."""
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    w_emu = int(width_pt * 12700)
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set("w", str(w_emu))
        solidFill = etree.SubElement(ln, qn("a:solidFill"))
        srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", hex_color.lstrip("#"))


def add_table(slide, rows_data, col_labels, row_labels,
              left, top, width, height):
    """
    Add a fully editable pptx table with:
      - Dark-blue header row
      - Light-blue baseline row (first data row)
      - Alternating white / very-light-grey rows
    """
    n_rows = len(rows_data)
    n_cols = len(col_labels)
    # Total rows in table = 1 header + n_rows data
    tbl_rows = 1 + n_rows
    tbl_cols = 1 + n_cols   # +1 for row-label column

    tbl = slide.shapes.add_table(
        tbl_rows, tbl_cols, left, top, width, height
    ).table

    # ── Column widths ──
    label_col_w = width * 0.16
    data_col_w  = (width - label_col_w) / n_cols
    tbl.columns[0].width = int(label_col_w)
    for c in range(1, tbl_cols):
        tbl.columns[c].width = int(data_col_w)

    # ── Row heights ──
    row_h = height // tbl_rows
    for r in range(tbl_rows):
        tbl.rows[r].height = int(row_h)

    def set_cell(r, c, text, bold=False, font_size=9,
                 fg=C_TEXT, align=PP_ALIGN.CENTER, bg=None, border=True):
        cell = tbl.cell(r, c)
        if bg:
            set_cell_color(cell, bg)
        if border:
            set_cell_border(cell, C_BORDER)
        tf = cell.text_frame
        tf.word_wrap = False
        p  = tf.paragraphs[0]
        p.alignment = align
        # Clear existing runs
        for run in p.runs:
            p._p.remove(run._r)
        run = p.add_run()
        run.text = str(text)
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = rgb(fg)

    # ── Header row (row 0) ──
    set_cell(0, 0, "", bg=C_HDR_BG, bold=True, fg=C_HDR_TEXT, font_size=9)
    for c, lbl in enumerate(col_labels):
        set_cell(0, c + 1, lbl, bg=C_HDR_BG, bold=True, fg=C_HDR_TEXT,
                 font_size=8)

    # ── Data rows ──
    for r, (row_label, row_vals) in enumerate(zip(row_labels, rows_data)):
        is_baseline = (r == 0)
        bg   = C_BL_ROW if is_baseline else (C_ROW_A if r % 2 == 0 else C_ROW_B)
        fg   = C_BLUE   if is_baseline else C_TEXT
        bold = is_baseline

        # Row label
        set_cell(r + 1, 0, row_label, bg=bg, bold=bold,
                 fg=fg, align=PP_ALIGN.LEFT, font_size=9)
        # Data cells
        for c, val in enumerate(row_vals):
            set_cell(r + 1, c + 1, val, bg=bg, bold=bold, fg=fg, font_size=9)


# ── Matplotlib plot (same style as standalone) ────────────────────────────────

def make_plot_fig(df, col, ylabel, baseline_val=None, tested_col="acc_is_test"):
    sub = df[df[col].notna()].copy()
    if sub.empty:
        return None

    BG       = "#ffffff"
    BLUE     = C_BLUE
    GRID     = C_GRID
    FRAME    = C_FRAME
    SUBTEXT  = C_SUBTEXT
    GREY     = C_GREY

    fig, ax = plt.subplots(figsize=(10.5, 4.8))
    fig.patch.set_facecolor(BG)
    ax.set_facecolor(BG)

    tested   = sub[sub[tested_col] == True]
    untested = sub[sub[tested_col] == False]

    ax.plot(sub["wf"], sub[col], "-", color=BLUE, linewidth=2.2, zorder=2, alpha=0.9)
    if not tested.empty:
        ax.plot(tested["wf"], tested[col], "o",
                color=BLUE, markersize=9, zorder=3, label="test result")
    if not untested.empty:
        ax.plot(untested["wf"], untested[col], "o",
                color=BLUE, markersize=9, zorder=3,
                markerfacecolor=BG, markeredgewidth=2.2,
                label="val acc epoch 39 (crashed — f=2 too coarse)")

    for _, row in sub.iterrows():
        val = row[col]
        lbl = f"{val:.3f}%" if col == "acc_%" else f"{val:.5f}"
        ax.annotate(lbl, xy=(row["wf"], val), xytext=(0, 11),
                    textcoords="offset points", ha="center",
                    fontsize=9, color=BLUE, fontweight="bold")

    if baseline_val is not None:
        ax.axhline(baseline_val, color=GREY, linestyle="--", linewidth=1.3,
                   alpha=0.9, label=f"FP baseline  {baseline_val}", zorder=1)

    for sp in ax.spines.values():
        sp.set_visible(True)
        sp.set_color(FRAME)
        sp.set_linewidth(0.9)
    ax.tick_params(colors=SUBTEXT, labelsize=11, length=0)
    ax.set_axisbelow(True)
    ax.grid(True, axis="y", linestyle="-", linewidth=0.8, alpha=0.65, color=GRID)
    ax.grid(False, axis="x")
    ax.set_xlabel("Weight fractional bits (wf)", fontsize=12, color=SUBTEXT, labelpad=6)
    ax.set_ylabel(ylabel, fontsize=12, color=SUBTEXT, labelpad=6)
    ax.set_xticks(sorted(sub["wf"].unique()))
    ymin, ymax = ax.get_ylim()
    ax.set_ylim(ymin, ymax + (ymax - ymin) * 0.14)

    if baseline_val is not None or not untested.empty:
        ax.legend(fontsize=10, framealpha=0.9, loc="lower right",
                  edgecolor=FRAME, labelcolor=SUBTEXT)

    plt.tight_layout()
    return fig


# ── Individual slide builders ─────────────────────────────────────────────────

def slide_plot(prs, df, col, title, ylabel, baseline_val=None):
    slide = blank_slide(prs)
    add_stripe(slide)
    add_title_block(slide, title, PLOT_SUBTITLE)

    fig = make_plot_fig(df, col, ylabel, baseline_val=baseline_val)
    if fig is not None:
        plot_top  = HEADER_H + Inches(0.05)
        plot_left = STRIPE_W + MARGIN
        plot_w    = SLIDE_W - STRIPE_W - 2 * MARGIN
        plot_h    = SLIDE_H - plot_top - Inches(0.25)
        embed_plot_image(slide, fig, plot_left, plot_top, plot_w, plot_h)


def slide_overview_table(prs, df):
    rows_data  = [BASELINE] + df.to_dict("records")
    row_labels = ["Full prec."] + [f"wf = {r['wf']}" for r in df.to_dict("records")]
    col_labels = ["Accuracy (%)", "Remaining weights", "OvO macro AUC"]

    rows = []
    for r in rows_data:
        acc_str = fmt_acc(r["acc_%"])
        if not r["acc_is_test"]:
            acc_str += " *"
        rows.append([acc_str, fmt_wt(r["remaining_w"]), fmt_auc(r["test_auc_ovo"])])

    slide = blank_slide(prs)
    add_stripe(slide)
    add_title_block(slide, "Overall scan metrics",
                    "Test accuracy  ·  remaining weights  ·  OvO macro AUC")
    add_table(slide, rows, col_labels, row_labels,
              left  = STRIPE_W + MARGIN,
              top   = HEADER_H,
              width = SLIDE_W - STRIPE_W - 2 * MARGIN,
              height= SLIDE_H - HEADER_H - Inches(0.55))
    add_footnote(slide,
        "* wf=2: training ran to epoch 39 but weights collapsed to 0% by epoch 3 "
        "(2-bit quantization too coarse); accuracy shown is 10% (random chance for 10 classes)")


def slide_auc_table(prs, df):
    rows_data  = [BASELINE] + df.to_dict("records")
    row_labels = ["Full prec."] + [f"wf = {r['wf']}" for r in df.to_dict("records")]
    col_labels = CLASSES

    rows = []
    for r in rows_data:
        rows.append([fmt_auc(r.get(f"auc_{c}")) for c in CLASSES])

    slide = blank_slide(prs)
    add_stripe(slide)
    add_title_block(slide, "Per-class AUC (OvR)",
                    "Binary one-vs-rest AUC per jet class")
    add_table(slide, rows, col_labels, row_labels,
              left  = STRIPE_W + MARGIN * 0.5,
              top   = HEADER_H,
              width = SLIDE_W - STRIPE_W - MARGIN,
              height= SLIDE_H - HEADER_H - Inches(0.45))
    add_footnote(slide, "FP baseline per-class AUC computed from pred.root")


def slide_rejection_table(prs, df):
    rows_data  = [BASELINE] + df.to_dict("records")
    row_labels = ["Full prec."] + [f"wf = {r['wf']}" for r in df.to_dict("records")]
    col_labels = [f"{cls} @{int(TPR_TARGETS[cls]*100)}%"
                  for cls in SIGNAL_CLASSES]

    rows = []
    for r in rows_data:
        rows.append([fmt_rej(r.get(f"rej_{c}")) for c in SIGNAL_CLASSES])

    slide = blank_slide(prs)
    add_stripe(slide)
    add_title_block(slide, "QCD rejection rate",
                    "1 / FPR vs QCD at per-class signal efficiency target")
    add_table(slide, rows, col_labels, row_labels,
              left  = STRIPE_W + MARGIN * 0.5,
              top   = HEADER_H,
              width = SLIDE_W - STRIPE_W - MARGIN,
              height= SLIDE_H - HEADER_H - Inches(0.55))
    add_footnote(slide,
        "Rejection = 1 / false positive rate, where FPR = fraction of QCD background jets passing the "
        "signal-score threshold at the given signal efficiency. "
        "Full prec. (FP) = full-precision uncompressed ParT model, evaluated on the 20M jet test set.")


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Generate editable PowerPoint slide deck for ParT wf scan")
    parser.add_argument("--output", default="scan_slides.pptx")
    args = parser.parse_args()

    print(f"Fetching runs from W&B ({ENTITY}/{PROJECT})...\n")
    df = fetch_full_data()

    if df.empty:
        print("\nNo data found. Check: wandb login")
        return

    print(f"\n--- Summary ---")
    print(df[["wf", "run_name", "acc_%", "test_auc_ovo", "state"]].to_string(index=False))

    print(f"\nBuilding PowerPoint → {args.output}")
    prs = new_presentation()

    slide_plot(prs, df,
               col="acc_%", title="ParT weights compression — accuracy vs wf bits",
               ylabel="Accuracy (%)", baseline_val=BASELINE_ACC)
    print("  Slide 1: accuracy plot")

    slide_plot(prs, df,
               col="test_auc_ovo", title="ParT weights compression — OvO AUC vs wf bits",
               ylabel="Test OvO macro AUC", baseline_val=BASELINE["test_auc_ovo"])
    print("  Slide 2: AUC plot")

    slide_overview_table(prs, df)
    print("  Slide 3: overview table")

    slide_auc_table(prs, df)
    print("  Slide 4: per-class AUC table")

    slide_rejection_table(prs, df)
    print("  Slide 5: rejection rates table")

    prs.save(args.output)
    print(f"\nDone — {args.output}  (5 slides, fully editable)")


if __name__ == "__main__":
    main()
